#! python2
import feedparser
import pyowm
import requests
import datetime

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches


# ############### - News - ################
# Function to fetch the rss feed and return the parsed RSS
def parseRSS(rss_url):
    return feedparser.parse(rss_url)


# Function returns the rss feed headlines and returns them as a list
def getHeadlines(feed):
    headlines = []
    for newsitem in feed['items']:
        headlines.append([newsitem['title'], newsitem['description']])
    return headlines


# Function takes an RSS URL, gets it parsed, and then asks for information
def getNews(rss_url):
    # Error if empty URL
    if rss_url == '':
        return 0
    # Get a parsed object for news source
    feed = parseRSS(rss_url)
    # A list to hold all headlines
    allheadlines = []
    # Call getHeadlines() and combine the returned info into a list of tuples
    # headlines and summaries
    allheadlines.extend(getHeadlines(feed))
    if 3 < len(allheadlines):
        limit = 3
    else:
        limit = len(allheadlines)
    return allheadlines[0:limit]


def getLocalNews(url, limit=3, detail=False):
    baseRequest = requests.get(url)
    baseData = baseRequest.text
    baseSoup = BeautifulSoup(baseData, "html.parser")
    artDump = baseSoup.find("div", {"id": "primary"})
    artDump_art = artDump.find_all("h1", {"class": "entry-title"})
    returnArticles = []
    for i in range(0, limit):
        art = artDump_art[i].get_text() + '\n'
        subUrl = artDump_art[i].find("a", href=True)['href']
        subR = requests.get(subUrl)
        subData = subR.text
        subSoup = BeautifulSoup(subData, "html.parser")
        subText = subSoup.find("div", {"class": "entry-content-wrap"})
        subArt = subText.find_all("p")
        if detail:
            for art in subArt:
                art += '  ' + art.get_text() + '\n'
        returnArticles.append(art)
    return returnArticles


def getScores(rss_url):
    # Error if empty URL
    if rss_url == '':
        return 0
    # Get a parsed object for news source
    feed = parseRSS(rss_url)
    # A list to hold all headlines
    allScores = []
    returnScores = []
    allScores.extend(getHeadlines(feed))
    if 3 < len(allScores):
        limit = 3
    else:
        limit = len(allScores)
    for i in range(0, limit):
        returnScores.append(allScores[i][1])
    return returnScores


# ############### - Weather - ################
def getDailyForecast():
    owm = pyowm.OWM('3cda239163ec0153dc83f9ae1bf3c54e')
    fc = owm.daily_forecast('Augusta,us', limit=4)
    f = fc.get_forecast()
    return f.get_weathers()


def getForecast():
    owm = pyowm.OWM('3cda239163ec0153dc83f9ae1bf3c54e')
    fc = owm.three_hours_forecast('Augusta,us')
    f = fc.get_forecast()
    return f.get_weathers()


def getWeather(limit=8):
    weekDays = ['Monday', 'Tuesday', 'Wednesday',
                'Thursday', 'Friday', 'Saturday', 'Sunday']
    forecast = ''
    for weather in getForecast():
        day = weekDays[weather.get_reference_time(
            'date').weekday()][0:3].upper()
        hour = weather.get_reference_time('date').hour - 4
        if hour < 0:
            hour += 24
        minute = weather.get_reference_time('date').minute
        w = str(weather.get_status())
        forecast += str(day + ' ' + str(hour) + ':' + str(
            minute) + ' > ' + w + '\n')
    return forecast.split('\n')[0:limit]


# ############### - Trivia - ################
def getTrivia():
    url = 'http://en.wikipedia.org/wiki/Main_Page'
    r = requests.get(url)
    data = r.text
    soup = BeautifulSoup(data, "html.parser")
    # Wikipedia On This Day div
    otd = soup.find("div", {"id": "mp-otd"})
    # ul items w/in ptd
    otd_ul = otd.find_all("ul")
    otd_p = otd.find_all("p")
    returnTrivia = []
    for item in otd_ul[0].get_text().encode('utf-8').split('\n'):
        returnTrivia.append(item)
    for item in otd_p:
        returnTrivia.append(item.get_text())
    # remove empty entries
    returnTrivia = [x for x in returnTrivia if x != '']
    return returnTrivia


# ############### - Trivia - ################
def getSources(sources):
    for source in sources:
        print source


# ############### - Body - ################

# Prints everything to terminal. Mostly used for testing the above
# functions to ensure that everything works. Not used right now.
def printReport():
    # List of RSS feeds that we will fetch and combine
    newsurls = {
        # test/local URLs
        # 'bbc-world': 'http://127.0.0.1:8000/bbcworld.xml',
        # 'bbc-us': 'http://127.0.0.1:8000/bbcus.xml',
        # 'local': 'http://127.0.0.1:8000/local.html',
        # 'sports': 'http://127.0.0.1:8000/soccer.xml'
        # normal URLs
        'bbc-world': 'http://feeds.bbci.co.uk/news/world/rss.xml',
        'bbc-us': 'http://feeds.bbci.co.uk/news/world/us_and_canada/rss.xml',
        'local': 'http://wjbf.com/category/news/csra-news/',
        'sports': 'http://www.scorespro.com/rss2/live-soccer.xml'
    }
    # List of other sources used
    sources = [
        'www.openweathermap.org',
        'en.wikipedia.org'
    ]
    # For each item in our newsurls
    for name, url in newsurls.items():
        # If it's not empy add it to our sources list for citation
        if url:
            sources.append(url)
        # Pretty printing
        print '\n----------------\n' + name + '\n'
        # call getNews() on the url
        if name != 'local':
            if name != 'sports':
                print getNews(url)
            else:
                print getScores(url)
        else:
            print getLocalNews(url)

    # Pretty printing
    print '\n----------------\nWeather\n'
    # Call getWeather() and for each day returned, print it
    for day in getWeather():
        print day

    # Pretty printing
    print '\n----------------\nOn This Day\n'
    # Pull some trivia from wikipedia
    print getTrivia()

    # Pretty printing
    print('\n----------------\nSources\n')
    # Print our sources
    getSources(sources)


# The powerpoint makes. Makes slides and calls relevant function for
# each one in order to fill it with content. Could go way deeper with this
# and automatically insert images from google based on the slide content
# (can't see that going wrong...) or resizing text and adding formatting.
# But for now this is plenty.
# Probably a way to shrink the size of this and make it more object oriented,
# but I don't want to do that right now.
# Can feed it your name for the title defaults to mine (because I can)
def makePresentation(name='CPT Devens'):
    # set up presentation
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    # set up title slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    # fill title slide
    date = datetime.datetime.now()
    title.text = "Daily Brief: " + str(
        date.now().month) + '/' + str(
        date.now().day) + '/' + str(
        date.now().year)
    subtitle.text = name
    # add logo
    img_path = 'logo_clear.png'
    left = top = Inches(0.64)
    height = Inches(0.6)
    pic = slide.shapes.add_picture(img_path, left, top, height=height)

    # set up world news slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    # fill first slide
    title.text = "World News"
    # slide.placeholders[1].text = "This is a test"
    text_frame = slide.placeholders[1].text_frame
    paragraph_strs = []
    url = 'http://feeds.bbci.co.uk/news/world/rss.xml'
    for pair in getNews(url):
        paragraph_strs.append(pair[0] + '. ' + pair[1])
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = paragraph_strs[0]
    for para_str in paragraph_strs[1:]:
        p = text_frame.add_paragraph()
        p.text = para_str

    # set up us news slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    # fill first slide
    title.text = "US / Canada News"
    # slide.placeholders[1].text = "This is a test"
    text_frame = slide.placeholders[1].text_frame
    paragraph_strs = []
    url = 'http://feeds.bbci.co.uk/news/world/us_and_canada/rss.xml'
    for pair in getNews(url):
        paragraph_strs.append(pair[0] + '. ' + pair[1])
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = paragraph_strs[0]
    for para_str in paragraph_strs[1:]:
        p = text_frame.add_paragraph()
        p.text = para_str

    # set up local news slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    # fill first slide
    title.text = "CSRA News"
    # slide.placeholders[1].text = "This is a test"
    text_frame = slide.placeholders[1].text_frame
    url = 'http://wjbf.com/category/news/csra-news/'
    paragraph_strs = getLocalNews(url)
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = paragraph_strs[0]
    for para_str in paragraph_strs[1:]:
        p = text_frame.add_paragraph()
        p.text = para_str

    # set up sports slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    # fill first slide
    title.text = "Sports"
    # slide.placeholders[1].text = "This is a test"
    text_frame = slide.placeholders[1].text_frame
    url = 'http://www.scorespro.com/rss2/live-soccer.xml'
    paragraph_strs = getScores(url)
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = paragraph_strs[0]
    for para_str in paragraph_strs[1:]:
        p = text_frame.add_paragraph()
        p.text = para_str

    # set up weather slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    # fill first slide
    title.text = "Weather"
    # slide.placeholders[1].text = "This is a test"
    text_frame = slide.placeholders[1].text_frame
    paragraph_strs = getWeather()
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = paragraph_strs[0]
    for para_str in paragraph_strs[1:]:
        p = text_frame.add_paragraph()
        p.text = para_str

    # set up trivia slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    # fill first slide
    title.text = "Trivia"
    # slide.placeholders[1].text = "This is a test"
    text_frame = slide.placeholders[1].text_frame
    paragraph_strs = getTrivia()
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = paragraph_strs[0]
    for para_str in paragraph_strs[1:]:
        p = text_frame.add_paragraph()
        p.text = para_str

    sources = [
        'http://feeds.bbci.co.uk/news/world/rss.xml',
        'http://feeds.bbci.co.uk/news/world/us_and_canada/rss.xml',
        'http://wjbf.com/category/news/csra-news/',
        'http://www.scorespro.com/rss2/live-soccer.xml',
        'http://www.openweathermap.org',
        'http://en.wikipedia.org'
    ]
    # set up sources slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    # fill first slide
    title.text = "Sources"
    # slide.placeholders[1].text = "This is a test"
    text_frame = slide.placeholders[1].text_frame
    paragraph_strs = sources
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = paragraph_strs[0]
    for para_str in paragraph_strs[1:]:
        p = text_frame.add_paragraph()
        p.text = para_str

    dateStr = str(date.year)
    if date.month < 10:
        dateStr = dateStr + '0' + str(date.month) + str(date.day)
    else:
        dateStr = dateStr + str(date.month) + str(date.day)
    prs.save('dailyBrief' + dateStr + '.pptx')


# printReport()
makePresentation()
